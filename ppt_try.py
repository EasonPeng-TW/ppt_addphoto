from pptx import Presentation
from pptx.enum.shapes import MSO_CONNECTOR
from pptx.util import Pt
from pptx.dml.color import RGBColor
from pptx.dml.line import LineFormat
from pptx.shapes.connector import Connector

LAYOUT_INDEX = 1
left = Pt(20)
top = Pt(20)
height = Pt(100)

presentation = Presentation()
slide_layout = presentation.slide_layouts[LAYOUT_INDEX]
slide = presentation.slides.add_slide(slide_layout)

# Patch the Connector class.
def get_or_add_ln(self):
    return self._element.spPr.get_or_add_ln()
Connector.get_or_add_ln = get_or_add_ln

# Draw a black line.
connector = slide.shapes.add_connector(
    MSO_CONNECTOR.STRAIGHT, left, top, left, top + height
)
# connector.ln = connector.get_or_add_ln()
line = LineFormat(connector)
line.fill.solid()
line.fill.fore_color.rgb = RGBColor(255, 0, 0)

presentation.save('./output.pptx')