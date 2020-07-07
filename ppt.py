import pptx.util
from pptx import Presentation
from pptx.util import Inches
from pptx.enum.shapes import MSO_CONNECTOR
from pptx.dml.color import RGBColor
from pptx.dml.line import LineFormat
from pptx.shapes.connector import Connector
from pptx.util import Cm
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.dml import MSO_COLOR_TYPE

# def get_or_add_ln(self):
#     return self._element.spPr.get_or_add_ln()
# Connector.get_or_add_ln = get_or_add_ln

prs = Presentation()
prs.slide_width = Inches(13.3) #adjust slide size
prs.slide_height = Inches(7.67)
title_slide_layout = prs.slide_layouts[0]
slide = prs.slides.add_slide(title_slide_layout)
slide2 = prs.slides.add_slide(title_slide_layout)
slide3 = prs.slides.add_slide(title_slide_layout)
slide4 = prs.slides.add_slide(title_slide_layout)
slide5 = prs.slides.add_slide(title_slide_layout)
slide6 = prs.slides.add_slide(title_slide_layout)
slide7 = prs.slides.add_slide(title_slide_layout)
slide8 = prs.slides.add_slide(title_slide_layout)
slide9 = prs.slides.add_slide(title_slide_layout)
slide10 = prs.slides.add_slide(title_slide_layout)
slide11 = prs.slides.add_slide(title_slide_layout)
slide12 = prs.slides.add_slide(title_slide_layout)
slide13 = prs.slides.add_slide(title_slide_layout)
slide14 = prs.slides.add_slide(title_slide_layout)
title = slide.shapes.title
subtitle = slide.placeholders[1]



# left = top = width = height = Cm(1.0)
# slide.shapes.add_shape(
#     MSO_SHAPE.LINE_CALLOUT_1, left, top, width, height
# )

left = Inches(0.4)
top = Inches(1.4)
try:
    img_path = '2.png'
    pic = slide.shapes.add_picture(img_path, left, top, width=pptx.util.Cm(30.41), height=pptx.util.Cm(7.27))
    # pic = slide.shapes.add_picture(img_path, left, top, width=pptx.util.Cm(8.98), height=pptx.util.Cm(7.27)) #for nomal size
# shapes = slide.shapes
    connector = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Cm(31.3), Cm(6.97), Cm(2.3), Cm(6.97)) #connector_type, begin_x, begin_y, end_x, end_y
    line = LineFormat(connector)
    line.fill.solid()
    line.fill.fore_color.rgb = RGBColor(255, 0, 0)
# assert shape.fill.fore_color.type == MSO_COLOR_TYPE.SCHEME
# fill = shape.fill
# fill.solid()
# fill.fore_color.rgb = RGBColor(255, 0, 0)
    # # line = LineFormat(connector.get_or_add_ln())
    # # line.fill.fore_color.rgb = RGBColor(0, 0, 0)    
    # connector.ln = connector.get_or_add_ln()
    # line = LineFormat(connector)
    # line.fill.solid()
    # line.fill.fore_color.rgb = RGBColor(0, 0, 0)
except:
    print(img_path + " error")

try:
    img_path = '3.png'
    pic = slide2.shapes.add_picture(img_path, left, top, width=pptx.util.Cm(30.41), height=pptx.util.Cm(7.27))
except :
    print(img_path + " error")
try:
    img_path = '4.png'
    pic = slide3.shapes.add_picture(img_path, left, top, width=pptx.util.Cm(30.41), height=pptx.util.Cm(7.27))
except:
    print(img_path + " error")
try:
    img_path = '5.png'
    pic = slide4.shapes.add_picture(img_path, left, top, width=pptx.util.Cm(30.41), height=pptx.util.Cm(7.27))
except:
    print(img_path + " error")
try:
    img_path = '7.png'
    pic = slide5.shapes.add_picture(img_path, left, top, width=pptx.util.Cm(30.41), height=pptx.util.Cm(7.27))
except:
    print(img_path + " error")
try:
    img_path = '8.png'
    pic = slide6.shapes.add_picture(img_path, left, top, width=pptx.util.Cm(30.41), height=pptx.util.Cm(7.27))
except:
    print(img_path + " error")
try:
    img_path = '12.png'
    pic = slide7.shapes.add_picture(img_path, left, top, width=pptx.util.Cm(30.41), height=pptx.util.Cm(7.27))
except:
    print(img_path + " error")
try:
    img_path = '13.png'
    pic = slide8.shapes.add_picture(img_path, left, top, width=pptx.util.Cm(30.41), height=pptx.util.Cm(7.27))
except:
    print(img_path + " error")
try:
    img_path = '20.png'
    pic = slide9.shapes.add_picture(img_path, left, top, width=pptx.util.Cm(30.41), height=pptx.util.Cm(7.27))
except:
    print(img_path + " error")
try:
    img_path = '26.png'
    pic = slide10.shapes.add_picture(img_path, left, top, width=pptx.util.Cm(30.41), height=pptx.util.Cm(7.27))
except:
    print(img_path + " error")
try:
    img_path = '29.png'
    pic = slide11.shapes.add_picture(img_path, left, top, width=pptx.util.Cm(30.41), height=pptx.util.Cm(7.27))
except:
    print(img_path + " error")
try:
    img_path = '30.png'
    pic = slide12.shapes.add_picture(img_path, left, top, width=pptx.util.Cm(30.41), height=pptx.util.Cm(7.27))
except:
    print(img_path + " error")
try:
    img_path = '41.png'
    pic = slide13.shapes.add_picture(img_path, left, top, width=pptx.util.Cm(30.41), height=pptx.util.Cm(7.27))
except:
    print(img_path + " error")
try:
    img_path = '66.png'
    pic = slide14.shapes.add_picture(img_path, left, top, width=pptx.util.Cm(30.41), height=pptx.util.Cm(7.27))
except:
    print(img_path + " error")


left = Inches(0.4)
top = Inches(4.4)

try:
    img_path = '2_off.png'
    pic = slide.shapes.add_picture(img_path, left, top, width=pptx.util.Cm(30.41), height=pptx.util.Cm(7.27))
except:
    print(img_path + " error")

try:
    img_path = '3_off.png'
    pic = slide2.shapes.add_picture(img_path, left, top, width=pptx.util.Cm(30.41), height=pptx.util.Cm(7.27))
except :
    print(img_path + " error")
try:
    img_path = '4_off.png'
    pic = slide3.shapes.add_picture(img_path, left, top, width=pptx.util.Cm(30.41), height=pptx.util.Cm(7.27))
except:
    print(img_path + " error")
try:
    img_path = '5_off.png'
    pic = slide4.shapes.add_picture(img_path, left, top, width=pptx.util.Cm(30.41), height=pptx.util.Cm(7.27))
except:
    print(img_path + " error")
try:
    img_path = '7_off.png'
    pic = slide5.shapes.add_picture(img_path, left, top, width=pptx.util.Cm(30.41), height=pptx.util.Cm(7.27))
except:
    print(img_path + " error")
try:
    img_path = '8_off.png'
    pic = slide6.shapes.add_picture(img_path, left, top, width=pptx.util.Cm(30.41), height=pptx.util.Cm(7.27))
except:
    print(img_path + " error")
try:
    img_path = '12_off.png'
    pic = slide7.shapes.add_picture(img_path, left, top, width=pptx.util.Cm(30.41), height=pptx.util.Cm(7.27))
except:
    print(img_path + " error")
try:
    img_path = '13_off.png'
    pic = slide8.shapes.add_picture(img_path, left, top, width=pptx.util.Cm(30.41), height=pptx.util.Cm(7.27))
except:
    print(img_path + " error")
try:
    img_path = '20_off.png'
    pic = slide9.shapes.add_picture(img_path, left, top, width=pptx.util.Cm(30.41), height=pptx.util.Cm(7.27))
except:
    print(img_path + " error")
try:
    img_path = '26_off.png'
    pic = slide10.shapes.add_picture(img_path, left, top, width=pptx.util.Cm(30.41), height=pptx.util.Cm(7.27))
except:
    print(img_path + " error")
try:
    img_path = '29_off.png'
    pic = slide11.shapes.add_picture(img_path, left, top, width=pptx.util.Cm(30.41), height=pptx.util.Cm(7.27))
except:
    print(img_path + " error")
try:
    img_path = '30_off.png'
    pic = slide12.shapes.add_picture(img_path, left, top, width=pptx.util.Cm(30.41), height=pptx.util.Cm(7.27))
except:
    print(img_path + " error")
try:
    img_path = '41_off.png'
    pic = slide13.shapes.add_picture(img_path, left, top, width=pptx.util.Cm(30.41), height=pptx.util.Cm(7.27))
except:
    print(img_path + " error")
try:
    img_path = '66_off.png'
    pic = slide14.shapes.add_picture(img_path, left, top, width=pptx.util.Cm(30.41), height=pptx.util.Cm(7.27))
except:
    print(img_path + " error")
prs.save('test.pptx')